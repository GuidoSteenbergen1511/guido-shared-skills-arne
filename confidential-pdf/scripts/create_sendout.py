#!/usr/bin/env python3
"""
create_sendout.py - Create send-out versions of workshop PPTX presentations.

Full pipeline: analyse slides -> remove unwanted -> add banner -> export PDF -> watermark -> compress.

Usage:
    python3 create_sendout.py \
        --source "/path/to/deck.pptx" \
        --client "ClientName" \
        --programme-type multi-session

    python3 create_sendout.py \
        --source "/path/to/deck.pptx" \
        --client "ClientName" \
        --banner-text "ClientName Team Only - Do Not Share" \
        --programme-type single-workshop \
        --no-watermark

    python3 create_sendout.py \
        --source "/path/to/deck.pptx" \
        --client "ClientName" \
        --dry-run

Slide removal rules derived from 4+ real workshop decks.
Banner dimensions reverse-engineered from actual sent files.
"""

import argparse
import os
import re
import shutil
import subprocess
import sys
import tempfile
from datetime import date
from difflib import SequenceMatcher
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("ERROR: python-pptx is not installed. Run: pip install python-pptx")
    sys.exit(1)

# Import watermark functions from sibling script
SCRIPT_DIR = Path(__file__).parent
sys.path.insert(0, str(SCRIPT_DIR))
from watermark_pdf import watermark_pdf, compress_pdf

# ---------------------------------------------------------------------------
# Banner constants (reverse-engineered from actual sent files - do not change)
# ---------------------------------------------------------------------------
BANNER_CX = 4796726     # EMU (~5.25 inches)
BANNER_CY = 334439      # EMU (~0.37 inches)
BANNER_X  = 3697637     # EMU (centred on 12192000 EMU wide slide)
BANNER_Y  = 0           # EMU (top of slide)
BANNER_RED = RGBColor(0xFF, 0x00, 0x00)

# Similarity thresholds for duplicate detection
NEAR_DUPLICATE_THRESHOLD = 0.80
EXACT_DUPLICATE_THRESHOLD = 0.90

# Default watermark/footer text
DEFAULT_WATERMARK = "CONFIDENTIAL \u2014 DO NOT DISTRIBUTE"
DEFAULT_FOOTER = (
    "CONFIDENTIAL \u2014 For personal use only. "
    "Do not distribute or share without written permission."
)


# ---------------------------------------------------------------------------
# Slide analysis
# ---------------------------------------------------------------------------

def extract_text(slide):
    """Return full text content of a slide as a single string."""
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            parts.append(shape.text_frame.text)
    return " ".join(parts).strip()


def text_similarity(a, b):
    """Return 0-1 similarity ratio between two strings."""
    if not a and not b:
        return 1.0
    if not a or not b:
        return 0.0
    return SequenceMatcher(None, a, b).ratio()


def has_video(slide):
    """Return True if slide contains a video/media element."""
    for shape in slide.shapes:
        if hasattr(shape, "media_type"):
            return True
        if shape.shape_type in (3, 13, 14, 15, 16):
            return True
    return False


def categorize_slide(slide, programme_type, seen_texts):
    """
    Analyse a slide and return (category, should_remove, reason).
    seen_texts: list of text strings from previous kept slides for duplicate detection.
    """
    all_text = extract_text(slide)
    text_len = len(all_text)
    lower = all_text.lower()

    # Rule 1: Blank / transition slides
    if text_len < 10:
        return "blank", True, "Fewer than 10 characters of text"

    # Rule 2: Video-only slides
    if has_video(slide) and text_len < 20:
        return "video", True, "Video/media element with minimal text"

    # Rule 5: Welcome / greeting slides
    greeting_patterns = [
        "nice to meet you", "welcome!", "welkom!", "welkom,", "goed om je te zien",
    ]
    if any(p in lower for p in greeting_patterns) and text_len < 100:
        return "welcome", True, "Short welcome/greeting slide"

    # Rule 6: Quote-only philosophical slides
    educational_keywords = [
        "programme", "exercise", "agenda", "key learnings", "recap",
        "step", "workflow", "module", "session", "training", "capability",
        "model", "prompt", "tool", "|",
        "customize", "configure", "create", "upload", "instructions",
        "chatgpt", "claude", "gpt", "project", "pull",
    ]
    is_educational = any(kw in lower for kw in educational_keywords)

    if not is_educational:
        has_attribution = bool(re.search(r'[\u2014\u2013\-]{1,2}\s*[A-Z]', all_text))
        starts_with_quote = bool(re.match(r'^[\s]*[\u201c\u201d\u2018\u2019"\u0022]', all_text))
        lines = [l.strip() for l in all_text.split('\n') if l.strip()]
        if (starts_with_quote or has_attribution) and text_len < 300 and len(lines) <= 4:
            return "quote", True, "Quote-only slide with no educational content"

    # Rules 3+4: Duplicate detection
    for prev_text in seen_texts:
        sim = text_similarity(all_text, prev_text)
        if sim >= EXACT_DUPLICATE_THRESHOLD:
            return "exact_duplicate", True, f"Exact duplicate (similarity {sim:.0%})"
        if sim >= NEAR_DUPLICATE_THRESHOLD:
            return "near_duplicate", True, f"Near-duplicate (similarity {sim:.0%})"

    # Multi-session only rules
    if programme_type == "multi-session":
        # Rule 7: Tool exploration exercises
        if all_text.upper().lstrip().startswith("EXERCISE"):
            tool_exploration_patterns = [
                "let's try", "laten we proberen",
                "let's test", "let's explore",
                "exploring", "go to chatgpt", "open chatgpt",
                "go to claude", "open claude",
                "ga naar chatgpt", "open claude.ai",
            ]
            if any(p in lower for p in tool_exploration_patterns):
                return "exercise_exploration", True, "Tool exploration exercise (multi-session)"

        # Rule 8: In-session recap references
        recap_patterns = [
            "remember when you", "you already practiced",
            "remember this overview",
            "weet je nog dat jij", "je hebt al geoefend",
        ]
        if any(p in lower for p in recap_patterns):
            return "recap_reference", True, "In-session recap reference (multi-session)"

        # Rule 9: Lesson-learned series (1/6, 2/6 etc.)
        if re.search(r'\(\d+/\d+\)', all_text):
            return "lesson_detail", True, "Lesson series slide e.g. 1/6 (multi-session)"

    return "content", False, "Educational content slide (kept)"


# ---------------------------------------------------------------------------
# Banner insertion
# ---------------------------------------------------------------------------

def add_banner_shape(slide, text):
    """Add a red Do-Not-Share banner rectangle to a slide."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Emu(BANNER_X), Emu(BANNER_Y), Emu(BANNER_CX), Emu(BANNER_CY),
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = BANNER_RED
    line = shape.line
    line.color.rgb = BANNER_RED
    line.width = Emu(12700)  # 1pt
    tf = shape.text_frame
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.runs[0]
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run.font.bold = True
    run.font.size = Pt(9)


# ---------------------------------------------------------------------------
# PDF export via LibreOffice
# ---------------------------------------------------------------------------

def export_pdf_libreoffice(pptx_path, output_dir):
    """Export PPTX to PDF using LibreOffice headless. Returns PDF path or None."""
    soffice = shutil.which("soffice")
    if not soffice:
        # macOS app bundle path
        mac_path = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
        if os.path.exists(mac_path):
            soffice = mac_path
        else:
            print("[WARN] LibreOffice not found. Skipping PDF export.", file=sys.stderr)
            return None

    result = subprocess.run(
        [soffice, "--headless", "--convert-to", "pdf", str(pptx_path), "--outdir", str(output_dir)],
        capture_output=True, text=True, timeout=300,
    )

    pdf_path = Path(output_dir) / (Path(pptx_path).stem + ".pdf")
    if pdf_path.exists():
        print(f"[SAVED ] PDF exported: {pdf_path}")
        return pdf_path
    else:
        print(f"[WARN  ] PDF export failed: {result.stderr.strip()}", file=sys.stderr)
        return None


# ---------------------------------------------------------------------------
# Main processing
# ---------------------------------------------------------------------------

def process_presentation(source_path, output_dir, banner_text, programme_type,
                         dry_run=False, do_watermark=True, do_compress=True,
                         watermark_text=DEFAULT_WATERMARK, footer_text=DEFAULT_FOOTER):
    """Full pipeline: copy -> analyse -> remove -> banner -> save PPTX -> PDF -> watermark -> compress."""
    source_path = Path(source_path)
    if not source_path.exists():
        print(f"ERROR: Source file not found: {source_path}", file=sys.stderr)
        sys.exit(1)

    # Output naming
    today = date.today().strftime("%y%m%d")
    stem = source_path.stem
    out_name = f"{today} {stem} (Do Not Share)"
    out_dir = Path(output_dir) if output_dir else source_path.parent
    out_dir.mkdir(parents=True, exist_ok=True)
    pptx_out = out_dir / f"{out_name}.pptx"

    if not dry_run:
        shutil.copy2(source_path, pptx_out)
        prs = Presentation(str(pptx_out))
    else:
        prs = Presentation(str(source_path))

    slides = list(prs.slides)
    total_original = len(slides)
    print(f"\n[INFO] Loaded: {source_path.name} ({total_original} slides)")
    print(f"[INFO] Programme type: {programme_type}")
    print(f"[INFO] Banner text: {banner_text}\n")

    # Analyse all slides
    removal_plan = []
    seen_texts = []

    for idx, slide in enumerate(slides):
        text = extract_text(slide)
        is_first = (idx == 0)
        is_last = (idx == len(slides) - 1)

        if is_first:
            seen_texts.append(text)
            print(f"  Slide {idx+1:>3}: [KEEP   ] cover/title slide")
            continue
        if is_last:
            seen_texts.append(text)
            print(f"  Slide {idx+1:>3}: [KEEP   ] closing/thank-you slide")
            continue

        category, should_remove, reason = categorize_slide(slide, programme_type, seen_texts)

        if should_remove:
            removal_plan.append((idx, category, reason))
            excerpt = text[:60].replace('\n', ' ') if text else "(no text)"
            print(f"  Slide {idx+1:>3}: [REMOVE ] {category:<22} - {reason}")
            print(f"            Text preview: \"{excerpt}\"")
        else:
            seen_texts.append(text)
            excerpt = text[:50].replace('\n', ' ') if text else "(no text)"
            print(f"  Slide {idx+1:>3}: [KEEP   ] {excerpt[:60]}")

    # Summary
    categories_removed = {}
    for _, cat, _ in removal_plan:
        categories_removed[cat] = categories_removed.get(cat, 0) + 1

    print("\n" + "=" * 60)
    print("REMOVAL SUMMARY")
    print("=" * 60)
    print(f"  Original slides  : {total_original}")
    print(f"  Slides to remove : {len(removal_plan)}")
    print(f"  Slides remaining : {total_original - len(removal_plan)}")
    if total_original > 0:
        pct = len(removal_plan) / total_original * 100
        print(f"  Reduction        : {pct:.0f}%")
    print()
    for cat, count in sorted(categories_removed.items()):
        print(f"  {cat:<25} : {count}")
    print("=" * 60)

    if dry_run:
        print("\n[DRY RUN] No files were modified.")
        return None

    # Remove slides in reverse order
    xml_slides = prs.slides._sldIdLst
    for idx, _, _ in sorted(removal_plan, key=lambda x: x[0], reverse=True):
        xml_slides.remove(xml_slides[idx])

    # Add banner to every slide
    final_slides = list(prs.slides)
    for slide in final_slides:
        add_banner_shape(slide, banner_text)
    print(f"\n[INFO  ] Banner added to all {len(final_slides)} slides")

    # Save PPTX
    prs.save(str(pptx_out))
    print(f"[SAVED ] PPTX: {pptx_out}")

    # Export to PDF via LibreOffice
    print("\n[STEP  ] Exporting PDF via LibreOffice...")
    pdf_path = export_pdf_libreoffice(pptx_out, out_dir)
    if not pdf_path:
        print("[WARN  ] PDF export failed. PPTX is still available.")
        return pptx_out

    # Rename PDF to match PPTX name (LibreOffice uses stem of input)
    final_pdf = out_dir / f"{out_name}.pdf"
    if pdf_path != final_pdf:
        pdf_path.rename(final_pdf)
        pdf_path = final_pdf

    # Watermark the PDF
    if do_watermark:
        print("[STEP  ] Adding PDF watermark...")
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            tmp_path = tmp.name
        try:
            # Rename current PDF to temp, watermark into final location
            shutil.move(str(pdf_path), tmp_path)
            if do_compress:
                # Watermark to another temp, then compress to final
                with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp2:
                    tmp2_path = tmp2.name
                try:
                    watermark_pdf(tmp_path, tmp2_path, watermark_text, footer_text)
                    compress_pdf(tmp2_path, str(pdf_path))
                finally:
                    if os.path.exists(tmp2_path):
                        os.unlink(tmp2_path)
            else:
                watermark_pdf(tmp_path, str(pdf_path), watermark_text, footer_text)
        finally:
            if os.path.exists(tmp_path):
                os.unlink(tmp_path)
    elif do_compress:
        print("[STEP  ] Compressing PDF...")
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            tmp_path = tmp.name
        try:
            shutil.move(str(pdf_path), tmp_path)
            compress_pdf(tmp_path, str(pdf_path))
        finally:
            if os.path.exists(tmp_path):
                os.unlink(tmp_path)

    size_mb = os.path.getsize(str(pdf_path)) / (1024 * 1024)
    print(f"[SAVED ] PDF: {pdf_path} ({size_mb:.1f} MB)")
    print("\n[DONE  ] Send-out creation complete.")
    return pptx_out


def main():
    parser = argparse.ArgumentParser(
        description="Create send-out version of workshop PPTX",
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument("--source", required=True, help="Path to original PPTX file")
    parser.add_argument("--client", required=True, help="Client name for banner text")
    parser.add_argument("--output-dir", default=None, help="Output directory (default: same as source)")
    parser.add_argument("--banner-text", default=None, help="Override banner text")
    parser.add_argument(
        "--programme-type",
        choices=["multi-session", "single-workshop"],
        default="multi-session",
        help="Affects which slide categories are removed",
    )
    parser.add_argument("--dry-run", action="store_true", help="Analyse only, do not write files")
    parser.add_argument("--no-watermark", action="store_true", help="Skip PDF diagonal watermark")
    parser.add_argument("--no-compress", action="store_true", help="Skip Ghostscript compression")

    args = parser.parse_args()
    banner_text = args.banner_text or f"{args.client} Team Only - Do Not Share"

    process_presentation(
        source_path=args.source,
        output_dir=args.output_dir,
        banner_text=banner_text,
        programme_type=args.programme_type,
        dry_run=args.dry_run,
        do_watermark=not args.no_watermark,
        do_compress=not args.no_compress,
    )


if __name__ == "__main__":
    main()
